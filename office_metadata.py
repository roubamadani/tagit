from docx import Document
from pptx import Presentation
import tempfile
import os
from datetime import datetime

def extract_docx_metadata(file_path):
    """Extract metadata from Word documents with error handling"""
    try:
        doc = Document(file_path)
        props = doc.core_properties
        
        return {
            'CoreProperties': {
                'title': props.title or '',
                'author': props.author or '',
                'subject': props.subject or '',
                'keywords': props.keywords or '',
                'comments': props.comments or '',
                'last_modified_by': props.last_modified_by or '',
                'created': props.created.isoformat() if props.created else '',
                'modified': props.modified.isoformat() if props.modified else '',
                'category': props.category or '',
                'content_status': props.content_status or '',
                'identifier': props.identifier or '',
                'language': props.language or '',
                'revision': props.revision or '',
                'version': props.version or ''
            },
            'DocumentStats': {
                'paragraph_count': len(doc.paragraphs),
                'tables_count': len(doc.tables),
                'sections_count': len(doc.sections)
            }
        }
    except Exception as e:
        return {'Error': f'DOCX Metadata Error: {str(e)}'}

def update_docx_metadata(input_path, changes):
    """Update Word document metadata and return path to new file"""
    try:
        # Create temp output path
        output_path = f"modified_{os.path.basename(input_path)}"
        
        # Make a copy of the original file
        with open(input_path, 'rb') as src, open(output_path, 'wb') as dst:
            dst.write(src.read())
        
        # Open the copy and update metadata
        doc = Document(output_path)
        props = doc.core_properties
        
        # Apply changes
        for key, value in changes.items():
            if hasattr(props, key):
                setattr(props, key, value)
        
        # Force modification date update
        props.modified = datetime.now()
        
        # Save changes
        doc.save(output_path)
        
        # Verify the update
        updated_doc = Document(output_path)
        for key in changes:
            if str(getattr(updated_doc.core_properties, key, '')) != str(changes[key]):
                raise ValueError(f"Failed to update {key}")
        
        return output_path
        
    except Exception as e:
        # Clean up if failed
        if os.path.exists(output_path):
            os.remove(output_path)
        return {'Error': f'DOCX Update Error: {str(e)}'}

def extract_pptx_metadata(file_path):
    """Extract metadata from PowerPoint files with error handling"""
    try:
        prs = Presentation(file_path)
        props = prs.core_properties
        
        return {
            'CoreProperties': {
                'title': props.title or '',
                'author': props.author or '',
                'subject': props.subject or '',
                'keywords': props.keywords or '',
                'comments': props.comments or '',
                'last_modified_by': props.last_modified_by or '',
                'created': props.created.isoformat() if props.created else '',
                'modified': props.modified.isoformat() if props.modified else '',
                'category': props.category or '',
                'content_status': props.content_status or '',
                'identifier': props.identifier or '',
                'language': props.language or '',
                'revision': props.revision or '',
                'version': props.version or ''
            },
            'PresentationStats': {
                'slide_count': len(prs.slides),
                'notes_slide_count': len(prs.slides._sldIdLst),
                'master_slide_count': len(prs.slide_masters)
            }
        }
    except Exception as e:
        return {'Error': f'PPTX Metadata Error: {str(e)}'}

def update_pptx_metadata(input_path, changes):
    """Update PowerPoint metadata and return path to new file"""
    try:
        # Create temp output path
        output_path = f"modified_{os.path.basename(input_path)}"
        
        # Make a copy of the original file
        with open(input_path, 'rb') as src, open(output_path, 'wb') as dst:
            dst.write(src.read())
        
        # Open the copy and update metadata
        prs = Presentation(output_path)
        props = prs.core_properties
        
        # Apply changes
        for key, value in changes.items():
            if hasattr(props, key):
                setattr(props, key, value)
        
        # Force modification date update
        props.modified = datetime.now()
        
        # Save changes
        prs.save(output_path)
        
        # Verify the update
        updated_prs = Presentation(output_path)
        for key in changes:
            if str(getattr(updated_prs.core_properties, key, '')) != str(changes[key]):
                raise ValueError(f"Failed to update {key}")
        
        return output_path
        
    except Exception as e:
        # Clean up if failed
        if os.path.exists(output_path):
            os.remove(output_path)
        return {'Error': f'PPTX Update Error: {str(e)}'}